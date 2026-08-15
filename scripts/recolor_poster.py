"""Recolor pooja pg poster.pptx with a chosen palette; write to a new file.

Usage:
    python scripts/recolor_poster.py [palette] [output.pptx]

Palettes: 'enhanced' (botanical forest-green), 'seminar' (Oxford heritage navy+gold).

Rules:
- Do NOT alter any text (content, font, size, bold, italic, color, position).
- Do NOT alter any shape position or size.
- Only replace shape *fill* elements inside <p:spPr> and the slide background fill.
"""
import shutil
import sys

from lxml import etree
from pptx import Presentation

SRC = r"v:\temp\Shelock_Finder\pooja pg poster.pptx"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def solid(hex_color: str) -> str:
    return (
        f'<a:solidFill xmlns:a="{NS_A}">'
        f'<a:srgbClr val="{hex_color}"/>'
        f"</a:solidFill>"
    )


def grad_linear(top_hex: str, bot_hex: str, angle: int = 5400000) -> str:
    return (
        f'<a:gradFill xmlns:a="{NS_A}" flip="none" rotWithShape="1">'
        f"<a:gsLst>"
        f'<a:gs pos="0"><a:srgbClr val="{top_hex}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{bot_hex}"/></a:gs>'
        f"</a:gsLst>"
        f'<a:lin ang="{angle}" scaled="1"/>'
        f"<a:tileRect/>"
        f"</a:gradFill>"
    )


def grad_radial(center_hex: str, edge_hex: str) -> str:
    return (
        f'<a:gradFill xmlns:a="{NS_A}" flip="none" rotWithShape="1">'
        f"<a:gsLst>"
        f'<a:gs pos="0"><a:srgbClr val="{center_hex}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{edge_hex}"/></a:gs>'
        f"</a:gsLst>"
        f'<a:path path="circle">'
        f'<a:fillToRect l="50000" t="50000" r="50000" b="50000"/>'
        f"</a:path>"
        f"<a:tileRect/>"
        f"</a:gradFill>"
    )


HEADER_NAMES = (
    "TextBox 45",   # title bar
    "Rectangle 1",  # INTRODUCTION
    "Rectangle 47", # RESULTS AND DISCUSSION
    "Rectangle 51", # CONCLUSION
    "Rectangle 27", # REFERENCES
    "Rectangle 35", # Table 1 caption
    "TextBox 43",   # Fig.2 caption
    "TextBox 49",   # Fig.3 caption
    "Rectangle 50", # OBJECTIVES
    "Rectangle 38", # Fig.1 caption
    "Rectangle 46", # MATERIAL AND METHODS
)


def palette_enhanced() -> tuple[str, dict[str, str]]:
    """Botanical forest-green over warm ivory."""
    slide_bg = "F5EFE0"
    header = "1B4332"
    fills = {name: solid(header) for name in HEADER_NAMES}
    fills.update({
        "Rectangle 13": grad_linear("F0F6FB", "C7DBEA"),   # INTRO pale blue
        "Rectangle 5":  solid("FFF6DA"),                    # OBJECTIVES cream
        "Rectangle 15": grad_linear("FFF9E3", "F0DCA0"),    # M&M parchment
        "Rectangle 20": grad_radial("FFF6E8", "EFC6A0"),    # RESULTS peach
        "Rectangle 33": grad_linear("F1F7EA", "C8DDB2"),    # CONCLUSION sage
        "Rectangle 8":  grad_linear("F1F7EA", "C8DDB2"),    # REFERENCES sage
        "Rectangle 14": grad_radial("F1F7EA", "D8E5C4"),    # decorative sage
    })
    return slide_bg, fills


def palette_seminar() -> tuple[str, dict[str, str]]:
    """Oxford Heritage: deep navy + antique gold + warm ivory.

    Highest-authority academic look; ideal for a seminar banner viewed at ~1.5 m.
    Very high white-on-navy contrast; gold accent draws eye to Objectives + Results.
    """
    slide_bg = "FAF6EF"  # warm ivory paper
    header = "0B2E4F"    # Oxford deep navy
    fills = {name: solid(header) for name in HEADER_NAMES}
    fills.update({
        # Body panels: warm off-white variations, restrained on purpose
        "Rectangle 13": grad_linear("FFFFFF", "EDE6D6"),   # INTRO white → warm dove
        "Rectangle 5":  solid("FBEFD1"),                    # OBJECTIVES gold highlight
        "Rectangle 15": grad_linear("FFFFFF", "EFE6D3"),    # M&M warm neutral
        "Rectangle 20": grad_radial("FFFDF5", "D8B36A"),    # RESULTS antique-gold glow (hero)
        "Rectangle 33": grad_linear("FFFFFF", "EBE1CD"),    # CONCLUSION deeper cream
        "Rectangle 8":  grad_linear("FFFFFF", "EFEAE0"),    # REFERENCES near-white
        "Rectangle 14": grad_radial("FBEFD1", "E8D7A7"),    # decorative subtle gold halo
    })
    return slide_bg, fills


PALETTES = {
    "enhanced": palette_enhanced,
    "seminar":  palette_seminar,
}

DEFAULT_OUT = {
    "enhanced": r"v:\temp\Shelock_Finder\pooja pg poster - enhanced.pptx",
    "seminar":  r"v:\temp\Shelock_Finder\pooja pg poster - seminar.pptx",
}

FILL_TAGS = {"solidFill", "gradFill", "pattFill", "blipFill", "noFill", "grpFill"}


def replace_fill(spPr, new_fill_xml: str) -> None:
    for child in list(spPr):
        if etree.QName(child).localname in FILL_TAGS:
            spPr.remove(child)
    new_fill = etree.fromstring(new_fill_xml)
    insert_at = 0
    for idx, child in enumerate(spPr):
        if etree.QName(child).localname in {"xfrm", "prstGeom", "custGeom"}:
            insert_at = idx + 1
    spPr.insert(insert_at, new_fill)


def find_spPr(sh):
    for child in sh._element.iter():
        q = etree.QName(child)
        if q.localname == "spPr":
            return child
    return None


def main() -> None:
    palette_name = sys.argv[1] if len(sys.argv) > 1 else "enhanced"
    if palette_name not in PALETTES:
        raise SystemExit(f"Unknown palette {palette_name!r}. Choose: {list(PALETTES)}")
    dst = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_OUT[palette_name]
    slide_bg, fills = PALETTES[palette_name]()

    shutil.copy(SRC, dst)
    prs = Presentation(dst)
    slide = prs.slides[0]

    bgPr = slide.background._element.find(f".//{{{NS_P}}}bg/{{{NS_P}}}bgPr")
    if bgPr is None:
        bgPr = slide.background._element.find(f".//{{{NS_P}}}bgPr")
    if bgPr is not None:
        for child in list(bgPr):
            if etree.QName(child).localname in FILL_TAGS:
                bgPr.remove(child)
        bgPr.insert(0, etree.fromstring(solid(slide_bg)))
        if bgPr.find(f"{{{NS_A}}}effectLst") is None:
            bgPr.append(etree.fromstring(f'<a:effectLst xmlns:a="{NS_A}"/>'))
        print(f"Slide background -> #{slide_bg}")
    else:
        print("!! no bgPr found")

    seen: set[str] = set()
    for sh in slide.shapes:
        if sh.name not in fills:
            continue
        spPr = find_spPr(sh)
        if spPr is None:
            print(f"!! {sh.name}: no spPr, skipped")
            continue
        replace_fill(spPr, fills[sh.name])
        seen.add(sh.name)
        print(f"[{sh.shape_id}] {sh.name}: fill replaced")

    missing = set(fills) - seen
    if missing:
        print(f"!! shapes in palette not found on slide: {sorted(missing)}")

    prs.save(dst)
    print(f"\nPalette: {palette_name}\nSaved:   {dst}")


if __name__ == "__main__":
    main()
